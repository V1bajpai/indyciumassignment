from pptx import Presentation
from pptx.util import Inches
from wand.image import Image
import os

class IndyciumAssignment:

	def __init__(self, prs):
		self.prs = prs

	# First slide
	def firstSlide(self):
		firstPageLayout = prs.slide_layouts[0]
		firstSlide = prs.slides.add_slide(firstPageLayout)
		self.title = firstSlide.shapes.title.text = "Indycium's Assignment Presentation"
		pass

	# Slides of Images with watermark
	def midSlides(self):
		layout = prs.slide_layouts[1]
		for i in range(5):
			slide = prs.slides.add_slide(layout)
			self.title = slide.shapes.title.text = 'This is Image-' + str(i+1)
			self.sub = slide.placeholders[1].text = 'This is subtitle' + str(i+1)

			# test that image file exists or not
			assert os.path.exists('C:\Indycium_project' + '\image'+str(i+1)+'.jpg')

			with Image(filename ='image'+str(i+1)+'.jpg') as image:
				with Image(filename ='nike_black.png') as water:
					with image.clone() as watermark:
						watermark.watermark(water, 0.5, 10, 20)
						watermark.save(filename ='watermarkedImage'+str(i+1)+'.jpg')

						#test- New watermarked file is getting created or not
						assert os.path.exists('C:\Indycium_project' + '\watermarkedImage' + str(i + 1) + '.jpg')

			self.pic = slide.shapes.add_picture('C:\Indycium_project\watermarkedImage'+str(i+1)+'.jpg', Inches(1), Inches(2.5), width=Inches(3.125984), height=Inches(3.9055))
		pass

	# Last slide
	def lastSlide(self):
		lastPageLayout = prs.slide_layouts[0]
		lastSlide = prs.slides.add_slide(lastPageLayout)
		self.title = lastSlide.shapes.title.text = 'Thank you'
		self.sub = lastSlide.placeholders[1].text = 'Presented By: Vaibhav Bajpai'
		pass

	# Saving the ppt
	def savePpt(self):
		prs.save("indyciumAssignment.pptx")
		pass

	# TestCase: after saving .pptx file exists or not
	def test_savedFileExisted(self):
		assert os.path.exists('C:\Indycium_project\indyciumAssignment.pptx')


prs = Presentation()
indyciumAssignment = IndyciumAssignment(prs)
indyciumAssignment.firstSlide()
indyciumAssignment.midSlides()
indyciumAssignment.lastSlide()
indyciumAssignment.savePpt()
indyciumAssignment.test_savedFileExisted()